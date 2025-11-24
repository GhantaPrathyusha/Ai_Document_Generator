import os
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
from app.config import settings


def clean_line(line: str) -> str:
    """Remove bullet prefixes like •, -, *, etc."""
    return line.lstrip("•-* ").strip()


def export_to_pptx(text: str) -> str:
    export_dir = settings.EXPORT_DIR
    os.makedirs(export_dir, exist_ok=True)

    prs = Presentation()
    layout = prs.slide_layouts[1]

    # Prepare lines
    raw_lines = [l.strip() for l in text.split("\n") if l.strip()]
    cleaned_lines = [clean_line(l) for l in raw_lines]

    slides = []
    current_title = None
    current_bullets = []

    for line in cleaned_lines:

        # Detect slide header
        # e.g. "Slide 3: Conclusion"
        if line.lower().startswith("slide") and ":" in line:
            # Save previous slide
            if current_title or current_bullets:
                slides.append((current_title, current_bullets))
            
            # Extract title after colon
            current_title = line.split(":", 1)[1].strip()
            current_bullets = []
        else:
            # Normal bullet
            current_bullets.append(line)

    # Add last slide
    if current_title or current_bullets:
        slides.append((current_title, current_bullets))

    # If no slide structure detected (fallback)
    if not slides:
        return "ERROR: No slides parsed."

    # Build PPT slides
    for idx, (title_text, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(layout)

        # Title
        slide.shapes.title.text = title_text or f"Slide {idx + 1}"

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        first = True
        for bullet in bullets:
            if first:
                p = tf.paragraphs[0]
                p.text = bullet
                first = False
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    filename = f"output_{int(datetime.utcnow().timestamp())}.pptx"
    filepath = os.path.join(export_dir, filename)
    prs.save(filepath)

    return filepath





# import os
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from datetime import datetime
# from app.config import settings

# def export_to_pptx(text: str) -> str:
#     export_dir = settings.EXPORT_DIR
#     os.makedirs(export_dir, exist_ok=True)

#     prs = Presentation()
#     title_layout = prs.slide_layouts[1]

#     # Convert text into bullet points
#     lines = [l.strip() for l in text.split("\n") if l.strip()]

#     # Group into slide chunks (4–6 bullets per slide)
#     chunk_size = 5
#     slides = [lines[i:i + chunk_size] for i in range(0, len(lines), chunk_size)]

#     for idx, slide_lines in enumerate(slides):
#         slide = prs.slides.add_slide(title_layout)
#         title = slide.shapes.title
#         body = slide.placeholders[1]

#         title.text = f"Slide {idx + 1}"

#         tf = body.text_frame
#         tf.clear()

#         first = True
#         for line in slide_lines:
#             if first:
#                 p = tf.paragraphs[0]
#                 p.text = " " + line
#                 first = False
#             else:
#                 p = tf.add_paragraph()
#                 p.text = " " + line
#                 p.level = 0

#     # Save file
#     filename = f"output_{int(datetime.utcnow().timestamp())}.pptx"
#     filepath = os.path.join(export_dir, filename)
#     prs.save(filepath)

#     return filepath
